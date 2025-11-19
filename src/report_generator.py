"""
Advanced Report Generator
Multi-format report generation (PDF, Excel, PowerPoint) with executive summaries
"""

import pandas as pd
import numpy as np
from datetime import datetime
from typing import Dict, List, Any, Optional
import json
from io import BytesIO
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

# PDF Generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.graphics.charts.piecharts import Pie
    from reportlab.graphics.charts.linecharts import HorizontalLineChart
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# PowerPoint Generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Enhanced Excel with Charts
try:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, PieChart, LineChart, Reference
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils.dataframe import dataframe_to_rows
    OPENPYXL_CHARTS_AVAILABLE = True
except ImportError:
    OPENPYXL_CHARTS_AVAILABLE = False


class AdvancedReportGenerator:
    """Generate comprehensive reports in multiple formats"""

    REPORT_TYPES = {
        'executive_summary': {
            'name': 'Executive Summary Report',
            'description': 'High-level overview for C-suite',
            'sections': ['portfolio_overview', 'key_metrics', 'top_risks', 'recommendations']
        },
        'technical_deep_dive': {
            'name': 'Technical Deep Dive',
            'description': 'Detailed technical analysis',
            'sections': ['portfolio_overview', 'tech_health', 'dependencies', 'modernization']
        },
        'financial_analysis': {
            'name': 'Financial Analysis Report',
            'description': 'Cost breakdown and optimization opportunities',
            'sections': ['cost_overview', 'tco_breakdown', 'hidden_costs', 'optimization']
        },
        'risk_compliance': {
            'name': 'Risk & Compliance Report',
            'description': 'Risk assessment and compliance status',
            'sections': ['risk_overview', 'compliance_status', 'mitigation_plans', 'audit_trail']
        },
        'roadmap_strategy': {
            'name': 'Strategic Roadmap Report',
            'description': 'Prioritized action plan with timeline',
            'sections': ['roadmap_overview', 'phase_breakdown', 'quick_wins', 'resources']
        }
    }

    def __init__(self, df_applications: pd.DataFrame):
        """Initialize with application portfolio data"""
        self.df = df_applications.copy()
        self.report_data = {}
        self.generated_at = datetime.now()

    def generate_portfolio_overview(self) -> Dict[str, Any]:
        """Generate high-level portfolio overview"""

        total_apps = len(self.df)
        total_cost = self.df['Cost'].sum()
        avg_health = self.df['Tech Health'].mean()
        avg_value = self.df['Business Value'].mean()

        # Category breakdown
        category_stats = self.df.groupby('Category').agg({
            'Application Name': 'count',
            'Cost': 'sum',
            'Tech Health': 'mean',
            'Business Value': 'mean'
        }).round(2)

        # Health distribution
        health_distribution = {
            'Critical (1-3)': len(self.df[self.df['Tech Health'] <= 3]),
            'Poor (4-5)': len(self.df[(self.df['Tech Health'] > 3) & (self.df['Tech Health'] <= 5)]),
            'Fair (6-7)': len(self.df[(self.df['Tech Health'] > 5) & (self.df['Tech Health'] <= 7)]),
            'Good (8-9)': len(self.df[(self.df['Tech Health'] > 7) & (self.df['Tech Health'] <= 9)]),
            'Excellent (10)': len(self.df[self.df['Tech Health'] == 10])
        }

        # Value distribution
        value_distribution = {
            'Low Value (1-3)': len(self.df[self.df['Business Value'] <= 3]),
            'Medium Value (4-6)': len(self.df[(self.df['Business Value'] > 3) & (self.df['Business Value'] <= 6)]),
            'High Value (7-10)': len(self.df[self.df['Business Value'] > 6])
        }

        return {
            'summary': {
                'total_applications': total_apps,
                'total_annual_cost': total_cost,
                'average_health_score': round(avg_health, 2),
                'average_business_value': round(avg_value, 2),
                'report_date': self.generated_at.strftime('%Y-%m-%d %H:%M:%S')
            },
            'category_breakdown': category_stats.to_dict('index'),
            'health_distribution': health_distribution,
            'value_distribution': value_distribution
        }

    def generate_key_metrics(self) -> Dict[str, Any]:
        """Generate key performance metrics"""

        # TIME Framework distribution
        time_distribution = {}
        for idx, row in self.df.iterrows():
            health = row['Tech Health']
            value = row['Business Value']

            if health <= 5 and value <= 6:
                category = 'Eliminate'
            elif health > 5 and value <= 6:
                category = 'Tolerate'
            elif health <= 5 and value > 6:
                category = 'Migrate'
            else:
                category = 'Invest'

            time_distribution[category] = time_distribution.get(category, 0) + 1

        # Cost efficiency
        high_cost_low_value = len(self.df[(self.df['Cost'] > self.df['Cost'].median()) &
                                           (self.df['Business Value'] <= 5)])

        low_cost_high_value = len(self.df[(self.df['Cost'] < self.df['Cost'].median()) &
                                           (self.df['Business Value'] >= 7)])

        # Technical debt estimate (apps with health < 5)
        technical_debt_apps = self.df[self.df['Tech Health'] < 5]
        technical_debt_cost = technical_debt_apps['Cost'].sum()

        return {
            'time_framework': time_distribution,
            'efficiency_metrics': {
                'high_cost_low_value_apps': high_cost_low_value,
                'low_cost_high_value_apps': low_cost_high_value,
                'cost_efficiency_ratio': round(low_cost_high_value / max(high_cost_low_value, 1), 2)
            },
            'technical_debt': {
                'affected_applications': len(technical_debt_apps),
                'annual_cost_at_risk': technical_debt_cost,
                'percentage_of_portfolio': round(len(technical_debt_apps) / len(self.df) * 100, 1)
            }
        }

    def generate_top_risks(self, limit: int = 10) -> List[Dict[str, Any]]:
        """Identify top risk applications"""

        risks = []

        for idx, row in self.df.iterrows():
            app_name = row['Application Name']
            health = row['Tech Health']
            value = row['Business Value']
            cost = row['Cost']

            # Calculate risk score
            health_risk = (10 - health) * 10
            criticality = value * 10
            cost_exposure = min(100, (cost / 500000) * 100)

            risk_score = (health_risk * 0.4 + criticality * 0.35 + cost_exposure * 0.25)

            # Risk factors
            risk_factors = []
            if health <= 3:
                risk_factors.append('Critical technical health')
            if value >= 8:
                risk_factors.append('Mission-critical application')
            if cost > 200000:
                risk_factors.append('High cost exposure')

            risks.append({
                'app_name': app_name,
                'risk_score': round(risk_score, 1),
                'health': health,
                'business_value': value,
                'annual_cost': cost,
                'risk_factors': risk_factors
            })

        # Sort by risk score
        risks.sort(key=lambda x: x['risk_score'], reverse=True)

        return risks[:limit]

    def generate_recommendations(self) -> List[Dict[str, Any]]:
        """Generate actionable recommendations"""

        recommendations = []

        # 1. Retirement candidates
        retire_candidates = self.df[(self.df['Tech Health'] <= 3) & (self.df['Business Value'] <= 4)]
        if len(retire_candidates) > 0:
            recommendations.append({
                'priority': 'high',
                'category': 'Rationalization',
                'action': 'Retire low-value, unhealthy applications',
                'impact': f'{len(retire_candidates)} applications identified',
                'estimated_savings': retire_candidates['Cost'].sum(),
                'timeline': '6-12 months',
                'details': f'Applications: {", ".join(retire_candidates["Application Name"].head(5).tolist())}'
            })

        # 2. Modernization priorities
        modernize_candidates = self.df[(self.df['Tech Health'] <= 5) & (self.df['Business Value'] >= 7)]
        if len(modernize_candidates) > 0:
            recommendations.append({
                'priority': 'urgent',
                'category': 'Modernization',
                'action': 'Modernize critical applications with poor health',
                'impact': f'{len(modernize_candidates)} critical applications at risk',
                'estimated_savings': modernize_candidates['Cost'].sum() * 0.2,
                'timeline': '3-6 months',
                'details': f'Applications: {", ".join(modernize_candidates["Application Name"].head(5).tolist())}'
            })

        # 3. Cost optimization
        high_cost_apps = self.df[self.df['Cost'] > self.df['Cost'].quantile(0.75)]
        if len(high_cost_apps) > 0:
            recommendations.append({
                'priority': 'medium',
                'category': 'Cost Optimization',
                'action': 'Review and optimize high-cost applications',
                'impact': f'Top 25% of applications by cost',
                'estimated_savings': high_cost_apps['Cost'].sum() * 0.15,
                'timeline': '3-9 months',
                'details': f'{len(high_cost_apps)} applications consuming {(high_cost_apps["Cost"].sum() / self.df["Cost"].sum() * 100):.1f}% of budget'
            })

        # 4. Consolidation opportunities
        category_counts = self.df['Category'].value_counts()
        overlapping_categories = category_counts[category_counts > 5]
        if len(overlapping_categories) > 0:
            recommendations.append({
                'priority': 'medium',
                'category': 'Consolidation',
                'action': 'Consolidate redundant applications in overlapping categories',
                'impact': f'{len(overlapping_categories)} categories with 5+ applications',
                'estimated_savings': self.df['Cost'].sum() * 0.10,
                'timeline': '6-18 months',
                'details': f'Categories: {", ".join(overlapping_categories.head(3).index.tolist())}'
            })

        return recommendations

    def generate_cost_breakdown(self) -> Dict[str, Any]:
        """Generate detailed cost breakdown"""

        total_cost = self.df['Cost'].sum()

        # By category
        category_costs = self.df.groupby('Category')['Cost'].sum().sort_values(ascending=False)

        # By health tier
        health_tiers = {
            'Critical (1-3)': self.df[self.df['Tech Health'] <= 3]['Cost'].sum(),
            'Poor (4-5)': self.df[(self.df['Tech Health'] > 3) & (self.df['Tech Health'] <= 5)]['Cost'].sum(),
            'Fair (6-7)': self.df[(self.df['Tech Health'] > 5) & (self.df['Tech Health'] <= 7)]['Cost'].sum(),
            'Good (8-9)': self.df[(self.df['Tech Health'] > 7) & (self.df['Tech Health'] <= 9)]['Cost'].sum(),
            'Excellent (10)': self.df[self.df['Tech Health'] == 10]['Cost'].sum()
        }

        # By value tier
        value_tiers = {
            'Low Value (1-3)': self.df[self.df['Business Value'] <= 3]['Cost'].sum(),
            'Medium Value (4-6)': self.df[(self.df['Business Value'] > 3) & (self.df['Business Value'] <= 6)]['Cost'].sum(),
            'High Value (7-10)': self.df[self.df['Business Value'] > 6]['Cost'].sum()
        }

        # Top 10 most expensive
        top_10_expensive = self.df.nlargest(10, 'Cost')[['Application Name', 'Cost', 'Tech Health', 'Business Value']].to_dict('records')

        return {
            'total_annual_cost': total_cost,
            'cost_by_category': category_costs.to_dict(),
            'cost_by_health': health_tiers,
            'cost_by_value': value_tiers,
            'top_10_expensive': top_10_expensive,
            'average_cost_per_app': round(total_cost / len(self.df), 2)
        }

    def generate_executive_summary_report(self) -> Dict[str, Any]:
        """Generate complete executive summary report"""

        return {
            'report_type': 'executive_summary',
            'report_title': 'Application Portfolio - Executive Summary',
            'generated_at': self.generated_at.isoformat(),
            'sections': {
                'portfolio_overview': self.generate_portfolio_overview(),
                'key_metrics': self.generate_key_metrics(),
                'top_risks': self.generate_top_risks(limit=10),
                'recommendations': self.generate_recommendations()
            }
        }

    def generate_financial_analysis_report(self) -> Dict[str, Any]:
        """Generate financial analysis report"""

        return {
            'report_type': 'financial_analysis',
            'report_title': 'Application Portfolio - Financial Analysis',
            'generated_at': self.generated_at.isoformat(),
            'sections': {
                'portfolio_overview': self.generate_portfolio_overview(),
                'cost_breakdown': self.generate_cost_breakdown(),
                'key_metrics': self.generate_key_metrics(),
                'recommendations': [r for r in self.generate_recommendations()
                                   if r['category'] in ['Cost Optimization', 'Rationalization']]
            }
        }

    def generate_technical_report(self) -> Dict[str, Any]:
        """Generate technical deep dive report"""

        # Technical health analysis
        health_stats = self.df.groupby('Category').agg({
            'Tech Health': ['mean', 'min', 'max', 'count']
        }).round(2)

        # Apps needing attention
        critical_apps = self.df[self.df['Tech Health'] <= 3][
            ['Application Name', 'Tech Health', 'Business Value', 'Cost', 'Category']
        ].to_dict('records')

        poor_apps = self.df[(self.df['Tech Health'] > 3) & (self.df['Tech Health'] <= 5)][
            ['Application Name', 'Tech Health', 'Business Value', 'Cost', 'Category']
        ].to_dict('records')

        return {
            'report_type': 'technical_deep_dive',
            'report_title': 'Application Portfolio - Technical Analysis',
            'generated_at': self.generated_at.isoformat(),
            'sections': {
                'portfolio_overview': self.generate_portfolio_overview(),
                'health_statistics': health_stats.to_dict(),
                'critical_applications': critical_apps,
                'poor_health_applications': poor_apps,
                'recommendations': [r for r in self.generate_recommendations()
                                   if r['category'] in ['Modernization', 'Consolidation']]
            }
        }

    def export_to_json(self, report_data: Dict[str, Any]) -> str:
        """Export report to JSON format"""
        return json.dumps(report_data, indent=2, default=str)

    def export_to_excel(self, report_data: Dict[str, Any]) -> BytesIO:
        """Export report to Excel format with multiple sheets and embedded charts"""

        output = BytesIO()

        if OPENPYXL_CHARTS_AVAILABLE:
            # Use openpyxl directly for chart support
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet

            # 1. Summary Sheet with formatting
            ws_summary = wb.create_sheet("Summary", 0)
            summary = report_data['sections']['portfolio_overview']['summary']

            # Header styling
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True, size=12)

            # Write summary data with formatting
            ws_summary['A1'] = "Application Portfolio Summary"
            ws_summary['A1'].font = Font(bold=True, size=16)
            ws_summary.merge_cells('A1:B1')

            row = 3
            for key, value in summary.items():
                ws_summary[f'A{row}'] = key.replace('_', ' ').title()
                ws_summary[f'B{row}'] = value
                ws_summary[f'A{row}'].font = Font(bold=True)
                row += 1

            # 2. Category Breakdown with Bar Chart
            if 'category_breakdown' in report_data['sections']['portfolio_overview']:
                ws_cat = wb.create_sheet("Category Breakdown")
                df_categories = pd.DataFrame.from_dict(
                    report_data['sections']['portfolio_overview']['category_breakdown'],
                    orient='index'
                )

                # Write data
                for r_idx, row in enumerate(dataframe_to_rows(df_categories, index=True, header=True), 1):
                    for c_idx, value in enumerate(row, 1):
                        cell = ws_cat.cell(row=r_idx, column=c_idx, value=value)
                        if r_idx == 1:
                            cell.font = header_font
                            cell.fill = header_fill

                # Add bar chart for cost by category
                chart = BarChart()
                chart.title = "Cost by Category"
                chart.x_axis.title = "Category"
                chart.y_axis.title = "Annual Cost ($)"

                data = Reference(ws_cat, min_col=3, min_row=1, max_row=len(df_categories)+1)
                cats = Reference(ws_cat, min_col=1, min_row=2, max_row=len(df_categories)+1)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                chart.height = 10
                chart.width = 20
                ws_cat.add_chart(chart, f"F2")

            # 3. Health Distribution with Pie Chart
            ws_health = wb.create_sheet("Health Distribution")
            health_dist = report_data['sections']['portfolio_overview']['health_distribution']

            ws_health['A1'] = "Health Level"
            ws_health['B1'] = "Count"
            ws_health['A1'].font = header_font
            ws_health['B1'].font = header_font
            ws_health['A1'].fill = header_fill
            ws_health['B1'].fill = header_fill

            row = 2
            for level, count in health_dist.items():
                ws_health[f'A{row}'] = level
                ws_health[f'B{row}'] = count
                row += 1

            # Add pie chart
            pie = PieChart()
            labels = Reference(ws_health, min_col=1, min_row=2, max_row=row-1)
            data = Reference(ws_health, min_col=2, min_row=1, max_row=row-1)
            pie.add_data(data, titles_from_data=True)
            pie.set_categories(labels)
            pie.title = "Application Health Distribution"
            pie.height = 10
            pie.width = 15
            ws_health.add_chart(pie, "D2")

            # 4. Top Risks
            if 'top_risks' in report_data['sections']:
                ws_risks = wb.create_sheet("Top Risks")
                df_risks = pd.DataFrame(report_data['sections']['top_risks'])

                for r_idx, row in enumerate(dataframe_to_rows(df_risks, index=False, header=True), 1):
                    for c_idx, value in enumerate(row, 1):
                        cell = ws_risks.cell(row=r_idx, column=c_idx, value=value)
                        if r_idx == 1:
                            cell.font = header_font
                            cell.fill = header_fill

            # 5. Recommendations
            if 'recommendations' in report_data['sections']:
                ws_rec = wb.create_sheet("Recommendations")
                df_recommendations = pd.DataFrame(report_data['sections']['recommendations'])

                for r_idx, row in enumerate(dataframe_to_rows(df_recommendations, index=False, header=True), 1):
                    for c_idx, value in enumerate(row, 1):
                        cell = ws_rec.cell(row=r_idx, column=c_idx, value=value)
                        if r_idx == 1:
                            cell.font = header_font
                            cell.fill = header_fill

            # 6. Cost Breakdown with Charts
            if 'cost_breakdown' in report_data['sections']:
                cost_data = report_data['sections']['cost_breakdown']
                if 'top_10_expensive' in cost_data:
                    ws_expensive = wb.create_sheet("Top 10 Expensive")
                    df_expensive = pd.DataFrame(cost_data['top_10_expensive'])

                    for r_idx, row in enumerate(dataframe_to_rows(df_expensive, index=False, header=True), 1):
                        for c_idx, value in enumerate(row, 1):
                            cell = ws_expensive.cell(row=r_idx, column=c_idx, value=value)
                            if r_idx == 1:
                                cell.font = header_font
                                cell.fill = header_fill

                    # Add bar chart
                    chart = BarChart()
                    chart.title = "Top 10 Most Expensive Applications"
                    chart.x_axis.title = "Application"
                    chart.y_axis.title = "Annual Cost ($)"

                    data = Reference(ws_expensive, min_col=2, min_row=1, max_row=11)
                    cats = Reference(ws_expensive, min_col=1, min_row=2, max_row=11)
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                    chart.height = 12
                    chart.width = 20
                    ws_expensive.add_chart(chart, "F2")

            # 7. Full Portfolio
            ws_full = wb.create_sheet("Full Portfolio")
            for r_idx, row in enumerate(dataframe_to_rows(self.df, index=False, header=True), 1):
                for c_idx, value in enumerate(row, 1):
                    cell = ws_full.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill

            wb.save(output)
            output.seek(0)
            return output
        else:
            # Fallback to pandas ExcelWriter
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                summary = report_data['sections']['portfolio_overview']['summary']
                df_summary = pd.DataFrame([summary])
                df_summary.to_excel(writer, sheet_name='Summary', index=False)

                if 'category_breakdown' in report_data['sections']['portfolio_overview']:
                    df_categories = pd.DataFrame.from_dict(
                        report_data['sections']['portfolio_overview']['category_breakdown'],
                        orient='index'
                    )
                    df_categories.to_excel(writer, sheet_name='Category Breakdown')

                if 'top_risks' in report_data['sections']:
                    df_risks = pd.DataFrame(report_data['sections']['top_risks'])
                    df_risks.to_excel(writer, sheet_name='Top Risks', index=False)

                if 'recommendations' in report_data['sections']:
                    df_recommendations = pd.DataFrame(report_data['sections']['recommendations'])
                    df_recommendations.to_excel(writer, sheet_name='Recommendations', index=False)

                if 'cost_breakdown' in report_data['sections']:
                    cost_data = report_data['sections']['cost_breakdown']
                    if 'top_10_expensive' in cost_data:
                        df_expensive = pd.DataFrame(cost_data['top_10_expensive'])
                        df_expensive.to_excel(writer, sheet_name='Top 10 Expensive', index=False)

                self.df.to_excel(writer, sheet_name='Full Portfolio', index=False)

            output.seek(0)
            return output

    def export_to_csv(self, report_data: Dict[str, Any]) -> str:
        """Export report summary to CSV format"""

        # Create summary CSV
        rows = []

        # Add summary section
        if 'portfolio_overview' in report_data['sections']:
            summary = report_data['sections']['portfolio_overview']['summary']
            for key, value in summary.items():
                rows.append({'Section': 'Portfolio Overview', 'Metric': key, 'Value': value})

        # Add key metrics
        if 'key_metrics' in report_data['sections']:
            metrics = report_data['sections']['key_metrics']
            for section, data in metrics.items():
                if isinstance(data, dict):
                    for key, value in data.items():
                        rows.append({'Section': f'Metrics - {section}', 'Metric': key, 'Value': value})

        df_export = pd.DataFrame(rows)
        return df_export.to_csv(index=False)

    def generate_report(self, report_type: str = 'executive_summary') -> Dict[str, Any]:
        """Generate report of specified type"""

        if report_type not in self.REPORT_TYPES:
            return {'error': f'Unknown report type: {report_type}'}

        if report_type == 'executive_summary':
            return self.generate_executive_summary_report()
        elif report_type == 'financial_analysis':
            return self.generate_financial_analysis_report()
        elif report_type == 'technical_deep_dive':
            return self.generate_technical_report()
        else:
            # Default to executive summary structure
            return self.generate_executive_summary_report()

    def export_to_pdf(self, report_data: Dict[str, Any]) -> BytesIO:
        """Export report to PDF format with charts and tables"""

        if not REPORTLAB_AVAILABLE:
            raise ImportError("ReportLab is required for PDF export. Install with: pip install reportlab")

        output = BytesIO()
        doc = SimpleDocTemplate(output, pagesize=letter, topMargin=0.5*inch, bottomMargin=0.5*inch)

        # Container for PDF elements
        elements = []
        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1e3a8a'),
            spaceAfter=12,
            alignment=TA_CENTER
        )

        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2563eb'),
            spaceAfter=10,
            spaceBefore=15
        )

        # Title
        elements.append(Paragraph(report_data['report_title'], title_style))
        elements.append(Paragraph(f"Generated: {report_data['generated_at']}", styles['Normal']))
        elements.append(Spacer(1, 0.3*inch))

        # Executive Summary Section
        elements.append(Paragraph("Executive Summary", heading_style))

        summary = report_data['sections']['portfolio_overview']['summary']
        summary_data = [
            ['Metric', 'Value'],
            ['Total Applications', summary['total_applications']],
            ['Total Annual Cost', f"${summary['total_annual_cost']:,.2f}"],
            ['Average Health Score', f"{summary['average_health_score']}/10"],
            ['Average Business Value', f"{summary['average_business_value']}/10"]
        ]

        summary_table = Table(summary_data, colWidths=[3*inch, 2.5*inch])
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2563eb')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(summary_table)
        elements.append(Spacer(1, 0.3*inch))

        # Health Distribution Pie Chart
        elements.append(Paragraph("Application Health Distribution", heading_style))

        health_dist = report_data['sections']['portfolio_overview']['health_distribution']
        drawing = Drawing(400, 200)
        pie = Pie()
        pie.x = 150
        pie.y = 50
        pie.width = 100
        pie.height = 100
        pie.data = list(health_dist.values())
        pie.labels = list(health_dist.keys())
        pie.slices.strokeWidth = 0.5
        drawing.add(pie)
        elements.append(drawing)
        elements.append(Spacer(1, 0.2*inch))

        # Top Risks
        if 'top_risks' in report_data['sections']:
            elements.append(Paragraph("Top 10 Risk Applications", heading_style))

            risks = report_data['sections']['top_risks'][:5]  # Top 5 for PDF
            risk_data = [['Application', 'Risk Score', 'Health', 'Business Value', 'Annual Cost']]

            for risk in risks:
                risk_data.append([
                    risk['app_name'][:30],  # Truncate long names
                    f"{risk['risk_score']:.1f}",
                    f"{risk['health']}/10",
                    f"{risk['business_value']}/10",
                    f"${risk['annual_cost']:,.0f}"
                ])

            risk_table = Table(risk_data, colWidths=[2*inch, 1*inch, 0.8*inch, 1.2*inch, 1.2*inch])
            risk_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#dc2626')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            elements.append(risk_table)
            elements.append(Spacer(1, 0.3*inch))

        # Recommendations
        if 'recommendations' in report_data['sections']:
            elements.append(PageBreak())
            elements.append(Paragraph("Strategic Recommendations", heading_style))

            for idx, rec in enumerate(report_data['sections']['recommendations'], 1):
                rec_title = f"{idx}. {rec['action']}"
                elements.append(Paragraph(rec_title, styles['Heading3']))

                rec_details = f"""
                <b>Priority:</b> {rec['priority'].upper()}<br/>
                <b>Category:</b> {rec['category']}<br/>
                <b>Impact:</b> {rec['impact']}<br/>
                <b>Estimated Savings:</b> ${rec['estimated_savings']:,.2f}<br/>
                <b>Timeline:</b> {rec['timeline']}<br/>
                <b>Details:</b> {rec['details']}
                """
                elements.append(Paragraph(rec_details, styles['Normal']))
                elements.append(Spacer(1, 0.2*inch))

        # Build PDF
        doc.build(elements)
        output.seek(0)
        return output

    def export_to_powerpoint(self, report_data: Dict[str, Any]) -> BytesIO:
        """Export report to PowerPoint presentation"""

        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint export. Install with: pip install python-pptx")

        output = BytesIO()
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Define colors
        blue_color = RGBColor(37, 99, 235)
        dark_blue = RGBColor(30, 58, 138)

        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = report_data['report_title']
        subtitle.text = f"Generated: {report_data['generated_at']}\nApplication Rationalization Analysis"

        # Slide 2: Executive Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Executive Summary"

        tf = body_shape.text_frame
        summary = report_data['sections']['portfolio_overview']['summary']

        p = tf.paragraphs[0]
        p.text = f"Total Applications: {summary['total_applications']}"
        p.level = 0

        p = tf.add_paragraph()
        p.text = f"Total Annual Cost: ${summary['total_annual_cost']:,.2f}"
        p.level = 0

        p = tf.add_paragraph()
        p.text = f"Average Health Score: {summary['average_health_score']}/10"
        p.level = 0

        p = tf.add_paragraph()
        p.text = f"Average Business Value: {summary['average_business_value']}/10"
        p.level = 0

        # Slide 3: Health Distribution Chart
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        title = slide.shapes.title
        title.text = "Application Health Distribution"

        # Add pie chart
        chart_data = CategoryChartData()
        health_dist = report_data['sections']['portfolio_overview']['health_distribution']

        chart_data.categories = list(health_dist.keys())
        chart_data.add_series('Applications', list(health_dist.values()))

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_CHART_TYPE.PIE
        chart.legend.include_in_layout = False

        # Slide 4: Top Risks
        if 'top_risks' in report_data['sections']:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]

            title_shape.text = "Top Risk Applications"

            tf = body_shape.text_frame
            risks = report_data['sections']['top_risks'][:5]

            for risk in risks:
                p = tf.paragraphs[0] if risks.index(risk) == 0 else tf.add_paragraph()
                p.text = f"{risk['app_name']}: Risk Score {risk['risk_score']:.1f}"
                p.level = 0

                # Add details
                p = tf.add_paragraph()
                p.text = f"Health: {risk['health']}/10, Value: {risk['business_value']}/10, Cost: ${risk['annual_cost']:,.0f}"
                p.level = 1

        # Slide 5: Key Recommendations
        if 'recommendations' in report_data['sections']:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]

            title_shape.text = "Strategic Recommendations"

            tf = body_shape.text_frame
            recommendations = report_data['sections']['recommendations'][:4]  # Top 4

            for idx, rec in enumerate(recommendations):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = rec['action']
                p.level = 0

                p = tf.add_paragraph()
                p.text = f"{rec['category']} - {rec['timeline']}"
                p.level = 1

                p = tf.add_paragraph()
                p.text = f"Savings: ${rec['estimated_savings']:,.0f}"
                p.level = 1

        # Slide 6: Cost Breakdown
        if 'cost_breakdown' in report_data['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Cost Breakdown by Health Tier"

            cost_data = report_data['sections']['cost_breakdown']
            health_costs = cost_data['cost_by_health']

            # Add bar chart
            chart_data = CategoryChartData()
            chart_data.categories = list(health_costs.keys())
            chart_data.add_series('Annual Cost ($)', list(health_costs.values()))

            x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False

        # Slide 7: Next Steps
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = "Next Steps"

        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Review and validate recommendations with stakeholders"
        p.level = 0

        p = tf.add_paragraph()
        p.text = "Prioritize quick wins for immediate impact"
        p.level = 0

        p = tf.add_paragraph()
        p.text = "Develop detailed implementation roadmap"
        p.level = 0

        p = tf.add_paragraph()
        p.text = "Establish governance process for ongoing optimization"
        p.level = 0

        # Save to BytesIO
        prs.save(output)
        output.seek(0)
        return output

    def send_email_report(self, report_data: Dict[str, Any],
                          recipients: List[str],
                          smtp_config: Dict[str, Any],
                          format: str = 'pdf',
                          subject: str = None) -> Dict[str, Any]:
        """Send report via email with attachment"""

        try:
            # Generate report in specified format
            if format == 'pdf':
                attachment = self.export_to_pdf(report_data)
                filename = f"{report_data['report_type']}_report.pdf"
                mime_type = 'application/pdf'
            elif format == 'excel':
                attachment = self.export_to_excel(report_data)
                filename = f"{report_data['report_type']}_report.xlsx"
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif format == 'powerpoint':
                attachment = self.export_to_powerpoint(report_data)
                filename = f"{report_data['report_type']}_report.pptx"
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            else:
                return {'success': False, 'error': f'Unsupported format: {format}'}

            # Create email
            msg = MIMEMultipart()
            msg['From'] = smtp_config.get('from_email', 'noreply@company.com')
            msg['To'] = ', '.join(recipients)
            msg['Subject'] = subject or f"Application Portfolio Report - {report_data['report_type']}"

            # Email body
            summary = report_data['sections']['portfolio_overview']['summary']
            body = f"""
            <html>
            <body>
                <h2>Application Portfolio Report</h2>
                <p>Please find attached the {report_data['report_type']} report generated on {report_data['generated_at']}.</p>

                <h3>Key Highlights:</h3>
                <ul>
                    <li><b>Total Applications:</b> {summary['total_applications']}</li>
                    <li><b>Total Annual Cost:</b> ${summary['total_annual_cost']:,.2f}</li>
                    <li><b>Average Health Score:</b> {summary['average_health_score']}/10</li>
                    <li><b>Average Business Value:</b> {summary['average_business_value']}/10</li>
                </ul>

                <p>For detailed analysis, please review the attached report.</p>

                <p>Best regards,<br/>Application Rationalization Tool</p>
            </body>
            </html>
            """

            msg.attach(MIMEText(body, 'html'))

            # Attach report
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(attachment.read())
            encoders.encode_base64(part)
            part.add_header('Content-Disposition', f'attachment; filename={filename}')
            msg.attach(part)

            # Send email
            server = smtplib.SMTP(smtp_config.get('smtp_host', 'localhost'),
                                 smtp_config.get('smtp_port', 587))

            if smtp_config.get('use_tls', True):
                server.starttls()

            if smtp_config.get('username') and smtp_config.get('password'):
                server.login(smtp_config['username'], smtp_config['password'])

            server.send_message(msg)
            server.quit()

            return {
                'success': True,
                'message': f'Report sent to {len(recipients)} recipient(s)',
                'recipients': recipients,
                'format': format
            }

        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

    def get_available_reports(self) -> Dict[str, Any]:
        """Get list of available report types"""
        return self.REPORT_TYPES

    def get_export_capabilities(self) -> Dict[str, bool]:
        """Check which export formats are available"""
        return {
            'json': True,
            'csv': True,
            'excel': True,
            'excel_with_charts': OPENPYXL_CHARTS_AVAILABLE,
            'pdf': REPORTLAB_AVAILABLE,
            'powerpoint': PPTX_AVAILABLE,
            'email': True
        }
